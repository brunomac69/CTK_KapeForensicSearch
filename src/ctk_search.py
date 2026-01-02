import os
import re
import argparse
import sqlite3
import hashlib
import csv
import docx
import pandas as pd
from pptx import Presentation
import pypdf
import sys
import warnings

warnings.filterwarnings("ignore")

# ---------------------------------------------------------
# Utils
# ---------------------------------------------------------

def is_large_file(path, size_mb):
    try:
        return os.path.getsize(path) > (size_mb * 1024 * 1024)
    except:
        return True

def get_sha1(file_path):
    sha1 = hashlib.sha1()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(65536), b''):
                sha1.update(chunk)
        return sha1.hexdigest()
    except:
        return "ERROR_CALCULATING_HASH"

def get_regex_pattern(keywords_file):
    if not os.path.exists(keywords_file):
        print(f"ERRO: Ficheiro de keywords nao encontrado: {keywords_file}")
        sys.exit(1)
    with open(keywords_file, 'r', encoding='utf-8', errors='ignore') as f:
        keywords = [line.strip() for line in f if line.strip()]
    if not keywords:
        print("ERRO: O ficheiro de keywords esta vazio.")
        sys.exit(1)
    patterns = [re.escape(k).replace(r'\*', '.*') for k in keywords]
    return re.compile("|".join(patterns), re.IGNORECASE)

def search_binary(file_path, pattern):
    """Fallback para arquivos de texto simples ou binários (e-mails mbox)."""
    try:
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            try:
                content = raw_data.decode('utf-8')
            except:
                content = raw_data.decode('latin-1', errors='ignore')
            match = pattern.search(content)
            return match.group(0) if match else None
    except:
        return None

# ---------------------------------------------------------
# Core processing
# ---------------------------------------------------------

def process_file(path, pattern, fast=False, deep=False):
    ext = os.path.splitext(path)[1].lower().lstrip('.')

    if fast:
        if ext == 'pdf' and is_large_file(path, 10): return None
        if ext in ['sqlite', 'sqlite3', 'db', 'accdb', 'mdb'] and is_large_file(path, 20): return None

    try:
        # DOCX
        if ext == 'docx':
            doc = docx.Document(path)
            full_text = "\n".join(p.text for p in doc.paragraphs)
            m = pattern.search(full_text)
            return m.group(0) if m else None

        # XLSX / XLS
        elif ext in ['xlsx', 'xls']:
            xl = pd.ExcelFile(path)
            for sheet in xl.sheet_names:
                df = xl.parse(sheet).astype(str)
                for col in df.columns:
                    for val in df[col]:
                        m = pattern.search(val)
                        if m: return m.group(0)
            return None

        # PPTX
        elif ext == 'pptx':
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        m = pattern.search(shape.text)
                        if m: return m.group(0)
            return None

        # PDF
        elif ext == 'pdf':
            reader = pypdf.PdfReader(path)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    m = pattern.search(t)
                    if m: return m.group(0)
            return None

        # SQLITE
        elif ext in ['sqlite', 'sqlite3', 'db']:
            conn = sqlite3.connect(path)
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            for (table_name,) in cursor.fetchall():
                cursor.execute(f"SELECT * FROM {table_name}")
                for row in cursor.fetchall():
                    m = pattern.search(str(row))
                    if m:
                        conn.close()
                        return m.group(0)
            conn.close()
            return None

        # Fallback para arquivos sem extensão (mbox) ou modo DEEP
        if deep or not ext:
            return search_binary(path, pattern)
        return None

    except:
        return search_binary(path, pattern) if (deep or not ext) else None

# ---------------------------------------------------------
# Main
# ---------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="CTK_search - Office & Financial Keyword Scanner")
    parser.add_argument("-d", "--dir", required=True, help="Diretorio de evidencia")
    parser.add_argument("-f", "--file", required=True, help="Caminho para o keywords.txt")
    parser.add_argument("-o", "--output", choices=['txt', 'csv'], default='txt', help="Formato de saida")
    parser.add_argument("-out", "--outfile", required=True, help="Ficheiro de saida")
    parser.add_argument("--fast", action="store_true")
    parser.add_argument("--deep", action="store_true")

    args = parser.parse_args()

    # Limpeza para KAPE
    search_dir = os.path.abspath(args.dir.strip('"')).rstrip('\\')
    output_file = os.path.abspath(args.outfile.strip('"'))

    if args.fast and args.deep:
        print("ERRO: As flags --fast e --deep nao podem ser usadas em simultaneo.")
        sys.exit(1)

    pattern = get_regex_pattern(args.file)
    results = []

    print(f"[*] A pesquisar em: {search_dir}")

    for root, _, files in os.walk(search_dir):
        for file in files:
            file_path = os.path.join(root, file)
            match_word = process_file(file_path, pattern, args.fast, args.deep)
            
            if match_word:
                sha1 = get_sha1(file_path)
                ext = os.path.splitext(file)[1].upper().lstrip('.')
                if not ext: ext = "DATA"

                results.append({
                    'Status': 'MATCH',
                    'Keyword': match_word,
                    'Extension': ext,
                    'SHA1': sha1,
                    'Path': file_path
                })
                print(f"MATCH | {match_word:<12} | {ext:<5} | {file_path}")

    # Exportação
    if args.output == 'csv':
        fields = ['Status', 'Keyword', 'Extension', 'SHA1', 'Path']
        with open(output_file, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.DictWriter(f, fieldnames=fields)
            writer.writeheader()
            writer.writerows(results)
    else:
        with open(output_file, 'w', encoding='utf-8-sig') as f:
            for r in results:
                f.write(f"{r['Status']} | {r['Keyword']} | {r['Extension']} | {r['SHA1']} | {r['Path']}\n")

    print(f"[*] Resultados salvos em: {output_file}")

if __name__ == "__main__":
    main()

#depois de venv e requirements gera exe: pyinstaller --onefile --console --clean --name "ctk_search" search.py